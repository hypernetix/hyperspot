#!/usr/bin/env python3
"""
Generate test data files for XLSX and PPTX parser tests.

Run this script to create the test files:
    pip install openpyxl python-pptx
    python generate_xlsx_pptx_testdata.py
"""
import os
from pathlib import Path

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font
except ImportError:
    print("Please install openpyxl: pip install openpyxl")
    exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Please install python-pptx: pip install python-pptx")
    exit(1)


def create_xlsx_simple_data():
    """Create a simple XLSX file with basic data."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"

    # Add header row
    ws['A1'] = 'Name'
    ws['B1'] = 'Age'
    ws['C1'] = 'City'

    # Make header bold
    for cell in ws[1]:
        cell.font = Font(bold=True)

    # Add data rows
    data = [
        ('Alice', 30, 'New York'),
        ('Bob', 25, 'San Francisco'),
        ('Charlie', 35, 'Chicago'),
    ]

    for row_idx, (name, age, city) in enumerate(data, start=2):
        ws[f'A{row_idx}'] = name
        ws[f'B{row_idx}'] = age
        ws[f'C{row_idx}'] = city

    return wb


def create_xlsx_multi_sheet():
    """Create an XLSX file with multiple sheets."""
    wb = Workbook()

    # First sheet - Sales
    ws1 = wb.active
    ws1.title = "Sales"
    ws1['A1'] = 'Product'
    ws1['B1'] = 'Quantity'
    ws1['C1'] = 'Price'
    ws1.append(['Widget', 100, 9.99])
    ws1.append(['Gadget', 50, 19.99])

    # Second sheet - Inventory
    ws2 = wb.create_sheet("Inventory")
    ws2['A1'] = 'Item'
    ws2['B1'] = 'Stock'
    ws2.append(['Widget', 500])
    ws2.append(['Gadget', 200])

    # Third sheet - Summary
    ws3 = wb.create_sheet("Summary")
    ws3['A1'] = 'Total Products'
    ws3['B1'] = 2
    ws3['A2'] = 'Total Revenue'
    ws3['B2'] = 1999.00

    return wb


def create_xlsx_merged_cells():
    """Create an XLSX file with merged cells."""
    wb = Workbook()
    ws = wb.active
    ws.title = "MergedCells"

    # Merge cells A1:B1 for header
    ws.merge_cells('A1:B1')
    ws['A1'] = 'Merged Header'
    ws['A1'].font = Font(bold=True)

    # Add data rows
    ws['A2'] = 'Column A'
    ws['B2'] = 'Column B'
    ws['A3'] = 'Data 1'
    ws['B3'] = 'Data 2'
    ws['A4'] = 'Data 3'
    ws['B4'] = 'Data 4'

    return wb


def create_xlsx_formula_cells():
    """Create an XLSX file with formula cells."""
    wb = Workbook()
    ws = wb.active
    ws.title = "Formulas"

    # Add numeric cells
    ws['A1'] = 10
    ws['A2'] = 20

    # Add formula cell
    ws['A3'] = '=SUM(A1:A2)'

    # Add more examples
    ws['B1'] = 5
    ws['B2'] = 15
    ws['B3'] = '=AVERAGE(B1:B2)'

    return wb


def create_pptx_simple():
    """Create a simple PPTX with one slide."""
    prs = Presentation()

    # Add title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Welcome to the Presentation"
    subtitle.text = "This is a simple test presentation for file parser testing."

    return prs


def create_pptx_multi_slide():
    """Create a PPTX with multiple slides."""
    prs = Presentation()

    # Slide 1 - Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Multi-Slide Presentation"
    slide.placeholders[1].text = "Testing multiple slides"

    # Slide 2 - Content
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide 2: Content"
    body = slide.placeholders[1]
    body.text = "This is the content of slide 2."

    # Slide 3 - More content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide 3: More Content"
    body = slide.placeholders[1]
    body.text = "This is the content of slide 3."

    return prs


def create_pptx_with_table():
    """Create a PPTX with a table."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add a table
    rows, cols = 3, 3
    left = Inches(1)
    top = Inches(1)
    width = Inches(6)
    height = Inches(2)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header row
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Role"
    table.cell(0, 2).text = "Department"

    # Data rows
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "Engineer"
    table.cell(1, 2).text = "Engineering"

    table.cell(2, 0).text = "Bob"
    table.cell(2, 1).text = "Manager"
    table.cell(2, 2).text = "Operations"

    return prs


def create_pptx_with_list():
    """Create a PPTX with bullet points."""
    prs = Presentation()

    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = "Key Points"

    body = slide.placeholders[1]
    tf = body.text_frame

    tf.text = "First important point"

    p = tf.add_paragraph()
    p.text = "Second important point"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Third important point"
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Sub-point under third"
    p.level = 1

    return prs


def main():
    script_dir = Path(__file__).parent
    xlsx_dir = script_dir / "xlsx"
    pptx_dir = script_dir / "pptx"

    # Ensure directories exist
    xlsx_dir.mkdir(exist_ok=True)
    pptx_dir.mkdir(exist_ok=True)

    # Generate XLSX files
    print("Generating XLSX test files...")

    wb = create_xlsx_simple_data()
    wb.save(xlsx_dir / "simple_data.xlsx")
    print(f"  Created: {xlsx_dir / 'simple_data.xlsx'}")

    wb = create_xlsx_multi_sheet()
    wb.save(xlsx_dir / "multi_sheet.xlsx")
    print(f"  Created: {xlsx_dir / 'multi_sheet.xlsx'}")

    wb = create_xlsx_merged_cells()
    wb.save(xlsx_dir / "merged_cells.xlsx")
    print(f"  Created: {xlsx_dir / 'merged_cells.xlsx'}")

    wb = create_xlsx_formula_cells()
    wb.save(xlsx_dir / "formula_cells.xlsx")
    print(f"  Created: {xlsx_dir / 'formula_cells.xlsx'}")

    # Generate PPTX files
    print("Generating PPTX test files...")

    prs = create_pptx_simple()
    prs.save(pptx_dir / "simple_presentation.pptx")
    print(f"  Created: {pptx_dir / 'simple_presentation.pptx'}")

    prs = create_pptx_multi_slide()
    prs.save(pptx_dir / "multi_slide.pptx")
    print(f"  Created: {pptx_dir / 'multi_slide.pptx'}")

    prs = create_pptx_with_table()
    prs.save(pptx_dir / "presentation_with_table.pptx")
    print(f"  Created: {pptx_dir / 'presentation_with_table.pptx'}")

    prs = create_pptx_with_list()
    prs.save(pptx_dir / "presentation_with_list.pptx")
    print(f"  Created: {pptx_dir / 'presentation_with_list.pptx'}")

    print("\nDone! Test data files have been created.")


if __name__ == "__main__":
    main()
